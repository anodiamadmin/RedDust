import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 1. Initialise presentation with a modern 16:9 widescreen layout
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette Definitions (Modern Engineering Theme)
DARK_BG = RGBColor(0x11, 0x1E, 0x2E)      # Deep Navy Midnight
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xF9)     # Soft Light Gray
TEXT_DARK = RGBColor(0x1A, 0x25, 0x30)    # Near Black Charcoal
TEXT_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)   # Crisp White
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xFF)  # Electric Signal Blue
ACCENT_MUTED = RGBColor(0x6C, 0x7A, 0x89) # Muted Technical Slate

def set_slide_background(slide, color):
    """Fills the slide background with a solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_clean_textbox(slide, left, top, width, height):
    """Helper to add a textbox with zero margins for pixel-perfect alignment."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    return tf

# ==============================================================================
# SLIDE 1: Title Slide (Dark Theme)
# ==============================================================================
slide_layout = prs.slide_layouts[6] # Blank layout
slide1 = prs.slides.add_slide(slide_layout)
set_slide_background(slide1, DARK_BG)

tf1 = add_clean_textbox(slide1, Inches(1.0), Inches(2.5), Inches(11.333), Inches(3.0))

p1 = tf1.paragraphs[0]
p1.text = "From Airwaves to Audio Files"
p1.font.name = "Arial"
p1.font.size = Pt(48)
p1.font.bold = True
p1.font.color.rgb = TEXT_LIGHT
p1.space_after = Pt(12)

p2 = tf1.add_paragraph()
p2.text = "The Physics and Mathematics of Digital Audio Sampling"
p2.font.name = "Arial"
p2.font.size = Pt(24)
p2.font.color.rgb = ACCENT_BLUE
p2.space_after = Pt(24)

p3 = tf1.add_paragraph()
p3.text = "Understanding the Nyquist Theorem, ADC Pipelines, and Quantisation"
p3.font.name = "Arial"
p3.font.size = Pt(16)
p3.font.color.rgb = ACCENT_MUTED

# ==============================================================================
# DATA STRUCTURE: Content Slides (Light Theme)
# ==============================================================================
slides_data = [
    {
        "title": "The Core Challenge: Continuous vs. Discrete",
        "points": [
            "Physical Reality: True sound is an analog pressure wave, infinitely variable in both time and amplitude.",
            "Digital Reality: Computational architectures only process binary code (discrete 0s and 1s).",
            "The Bridge: Sampling acts as the mathematical and electrical link, slicing continuous waves into distinct numerical frames.",
            "Information Integrity: The process must capture enough detail to allow exact wave reconstruction without raw data bloat."
        ]
    },
    {
        "title": "The Architectural ADC Pipeline",
        "points": [
            "Acoustic Input: Physical sound pressure waves vibrate a microphone diaphragm, generating low-voltage electrical signals.",
            "Pre-Amplification: The low-voltage microphone signal is boosted to a clean, usable line-level hardware voltage.",
            "Anti-Aliasing Filter: A low-pass hardware filter violently strips away all frequencies above the chosen system limit.",
            "The ADC Chip: The core integrated circuit handles both Sampling (chopping time) and Quantisation (chopping amplitude).",
            "Storage & Packaging: The resulting structural binary stream is written directly into uncompressed .wav file containers."
        ]
    },
    {
        "title": "Step 1: Time Discretisation (Sampling)",
        "points": [
            "Definition: Measuring the exact instantaneous voltage of an analog signal at uniform, highly precise intervals of time.",
            "Sampling Interval (Ts): The fixed, unvarying time duration separating consecutive signal measurements.",
            "Sampling Frequency (fs): The total number of snapshots captured per single second, measured strictly in Hertz (Hz).",
            "Mathematical Transformation: Maps a continuous-time signal x(t) into a discrete sequence defined as x[n] = x(nTs)."
        ]
    },
    {
        "title": "The Nyquist-Shannon Sampling Theorem",
        "points": [
            "The Fundamental Rule: To perfectly reconstruct an analog signal, the sampling rate (fs) must be greater than twice the highest frequency component (fmax) contained within the signal.",
            "The Formula: fs >= 2 * fmax",
            "Nyquist Frequency: The absolute mathematical frequency ceiling of a digital system, equal to exactly half the sampling rate (fs / 2).",
            "Reconstruction: If this condition is met, the original wave can be rebuilt with zero mathematical loss or structural interpolation errors."
        ]
    },
    {
        "title": "The Penalty of Failure: Aliasing",
        "points": [
            "Under-sampling: Occurs when a sound wave changes faster than the digital system is configured to sample it (fs < 2 * fmax).",
            "Spectral Folding: The high-frequency wave reflects back across the Nyquist limit, impersonating an artificial low-frequency wave.",
            "Digital Distortion: Creates harsh, metallic noise that is permanently baked into the data stream and cannot be filtered out later.",
            "Prevention: The hardware anti-aliasing filter forces strict compliance by cutting off illegal high frequencies prior to conversion."
        ]
    },
    {
        "title": "Real-World Engineering: Why 44.1 kHz?",
        "points": [
            "Human Biological Limits: The healthy human auditory system can perceive frequency vibrations up to roughly 20,000 Hz (20 kHz).",
            "The Nyquist Baseline: According to the theorem, capturing human hearing requires an absolute minimum baseline of 40 kHz.",
            "The Transition Band: The extra 4,100 Hz provides critical breathing room for physical, real-world analog filters to roll off efficiently.",
            "Industry Standardization: Established during the early digital audio era to sync seamlessly with global video formats."
        ]
    },
    {
        "title": "Step 2: Amplitude Discretisation (Quantisation)",
        "points": [
            "The Axis Split: While sampling slices the horizontal axis (Time), quantisation slices the vertical axis (Voltage Level).",
            "Bit Depth (Resolution): The exact number of bits dedicated to expressing the numeric value of each individual sample point.",
            "Resolution Scale Levels:",
            "  - 16-Bit (CD Quality): Provides 2^16 = 65,536 distinct numerical amplitude steps.",
            "  - 24-Bit (Studio Master): Provides 2^24 = 16,777,216 distinct numerical amplitude steps.",
            "Quantisation Error: The tiny mathematical rounding offset between the true wave voltage and the nearest grid line. This error generates a permanent low-level noise floor."
        ]
    },
    {
        "title": "Pulse Code Modulation (PCM) & WAV Files",
        "points": [
            "PCM Generation: The sequential stream of quantised integers is encoded directly into a continuous digital bitstream.",
            "WAV Container Architecture:",
            "  - Header Chunk: Defines critical file identity parameters: sample rate, bit depth, and channel configurations (Mono/Stereo).",
            "  - Data Chunk: A massive, uncompressed sequential block of raw binary data representing audio amplitude over time.",
            "Data Rate Formula: Bit Rate = Sample Rate * Bit Depth * Number of Channels",
            "  - CD Calculation: 44,100 Hz * 16 bits * 2 channels = 1,411.2 kbps constant data throughput."
        ]
    },
    {
        "title": "Conclusion: Key Engineering Principles",
        "points": [
            "Time Domain Mastery: Elevating the sample rate extends the system's ability to record higher frequencies cleanly.",
            "Amplitude Domain Mastery: Increasing the system bit depth drops the noise floor and expands total operational dynamic range.",
            "Mathematical Perfection: When Nyquist criteria are fully satisfied, digital audio does not lose information—there are no jagged 'staircases' in the final output wave."
        ]
    }
]

# Generate each content slide sequentially using the structured data
for slide_data in slides_data:
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, LIGHT_BG)
    
    # Title Block
    tf_title = add_clean_textbox(slide, Inches(1.0), Inches(0.8), Inches(11.333), Inches(1.0))
    p_title = tf_title.paragraphs[0]
    p_title.text = slide_data["title"]
    p_title.font.name = "Arial"
    p_title.font.size = Pt(32)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_DARK
    
    # Body Content Box
    tf_body = add_clean_textbox(slide, Inches(1.0), Inches(2.2), Inches(11.333), Inches(4.5))
    
    for idx, point in enumerate(slide_data["points"]):
        p = tf_body.paragraphs[0] if idx == 0 else tf_body.add_paragraph()
        p.text = point
        p.font.name = "Arial"
        
        # Adjust typography styling dynamically for structural sub-bullets
        if point.startswith("  -"):
            p.font.size = Pt(16)
            p.font.color.rgb = ACCENT_MUTED
            p.space_after = Pt(8)
            p.level = 1
        else:
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_DARK
            p.space_after = Pt(14)
            p.level = 0

# 3. Save file locally
prs.save("Audio_Sampling_and_Nyquist_Rate.pptx")
